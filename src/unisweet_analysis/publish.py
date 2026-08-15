from __future__ import annotations

import json
from pathlib import Path
from typing import Any

import xlsxwriter
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


NAVY = "#0F0E9A"
BLUE = "#BAE1FF"
GREEN = "#BAFFC9"
RED = "#FFB3BA"
AMBER = "#FFE6A7"
PALE = "#F4F6FB"
INK = "#1C2434"
WHITE = "#FFFFFF"
GRID = "#D7DDEA"


def _signed(value: float | None, decimals: int = 0) -> str:
    return f"{float(value or 0):+,.{decimals}f}"


def _percent(value: float | None, decimals: int = 1) -> str:
    return f"{float(value or 0) * 100:.{decimals}f}%"


def _formats(workbook: xlsxwriter.Workbook) -> dict[str, Any]:
    return {
        "title": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 20, "bold": True, "font_color": WHITE, "bg_color": NAVY, "align": "left", "valign": "vcenter"}
        ),
        "subtitle": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 11, "font_color": WHITE, "bg_color": NAVY, "align": "left", "valign": "vcenter"}
        ),
        "section": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 10, "bold": True, "font_color": WHITE, "bg_color": NAVY, "align": "left", "valign": "vcenter"}
        ),
        "header": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 9, "bold": True, "font_color": WHITE, "bg_color": NAVY, "align": "center", "valign": "vcenter", "text_wrap": True, "border": 1, "border_color": WHITE}
        ),
        "body": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 9, "font_color": INK, "valign": "vcenter", "bottom": 1, "bottom_color": GRID}
        ),
        "body_wrap": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 9, "font_color": INK, "valign": "top", "text_wrap": True, "bottom": 1, "bottom_color": GRID}
        ),
        "amount": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 9, "font_color": INK, "num_format": "#,##0;[Red](#,##0);-", "align": "right", "bottom": 1, "bottom_color": GRID}
        ),
        "pct": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 9, "font_color": INK, "num_format": "0.0%;[Red](0.0%);-", "align": "right", "bottom": 1, "bottom_color": GRID}
        ),
        "bps": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 9, "font_color": INK, "num_format": '0 "bps";[Red](0 "bps");-', "align": "right", "bottom": 1, "bottom_color": GRID}
        ),
        "card_label": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 9, "bold": True, "font_color": WHITE, "bg_color": NAVY, "align": "center", "valign": "vcenter"}
        ),
        "card_blue": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 22, "bold": True, "font_color": INK, "bg_color": BLUE, "align": "center", "valign": "vcenter", "num_format": "#,##0", "border": 2, "border_color": NAVY}
        ),
        "card_red": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 22, "bold": True, "font_color": "#C00000", "bg_color": RED, "align": "center", "valign": "vcenter", "num_format": "#,##0;[Red](#,##0);-", "border": 2, "border_color": NAVY}
        ),
        "card_green": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 22, "bold": True, "font_color": INK, "bg_color": GREEN, "align": "center", "valign": "vcenter", "num_format": "#,##0", "border": 2, "border_color": NAVY}
        ),
        "card_pct": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 22, "bold": True, "font_color": INK, "bg_color": GREEN, "align": "center", "valign": "vcenter", "num_format": "0.0%", "border": 2, "border_color": NAVY}
        ),
        "insight_blue": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 8, "font_color": INK, "bg_color": BLUE, "valign": "top", "text_wrap": True, "border": 1, "border_color": NAVY}
        ),
        "insight_red": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 8, "font_color": INK, "bg_color": RED, "valign": "top", "text_wrap": True, "border": 1, "border_color": NAVY}
        ),
        "insight_green": workbook.add_format(
            {"font_name": "Segoe UI", "font_size": 8, "font_color": INK, "bg_color": GREEN, "valign": "top", "text_wrap": True, "border": 1, "border_color": NAVY}
        ),
        "pass": workbook.add_format({"font_name": "Segoe UI", "font_size": 9, "bold": True, "bg_color": GREEN, "font_color": INK}),
        "warn": workbook.add_format({"font_name": "Segoe UI", "font_size": 9, "bold": True, "bg_color": AMBER, "font_color": INK}),
        "fail": workbook.add_format({"font_name": "Segoe UI", "font_size": 9, "bold": True, "bg_color": RED, "font_color": INK}),
    }


def _write_title(sheet: Any, formats: dict[str, Any], title: str, subtitle: str, last_col: str = "O") -> None:
    sheet.merge_range(f"A1:{last_col}2", title, formats["title"])
    sheet.merge_range(f"A3:{last_col}3", subtitle, formats["subtitle"])
    sheet.set_row(0, 25)
    sheet.set_row(1, 25)
    sheet.set_row(2, 21)


def _write_table(
    sheet: Any,
    start_row: int,
    start_col: int,
    headers: list[str],
    rows: list[list[Any]],
    formats: dict[str, Any],
    column_formats: dict[int, str] | None = None,
) -> None:
    for column, header in enumerate(headers):
        sheet.write(start_row, start_col + column, header, formats["header"])
    for row_index, row in enumerate(rows, start=start_row + 1):
        for column, value in enumerate(row):
            style = formats[(column_formats or {}).get(column, "body")]
            sheet.write(row_index, start_col + column, value, style)


def publish_excel(pack: dict[str, Any], output_path: Path) -> None:
    if pack["metadata"]["publication_status"] != "READY":
        raise ValueError("Dashboard publishing requires publication_status=READY")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    workbook = xlsxwriter.Workbook(output_path)
    workbook.set_properties(
        {
            "title": "UniSweet Finance Analysis Dashboard",
            "subject": "Python-generated financial analysis",
            "author": "UniSweet Finance Analysis",
            "comments": "Generated from analysis_pack.json; no embedded source workbook links.",
        }
    )
    formats = _formats(workbook)
    dashboard = workbook.add_worksheet("Dashboard")
    sales = workbook.add_worksheet("Sales Analysis")
    leakage = workbook.add_worksheet("Discount Leakage")
    pnl_market = workbook.add_worksheet("P&L and Market")
    checks = workbook.add_worksheet("Checks and Sources")

    for sheet in [dashboard, sales, leakage, pnl_market, checks]:
        sheet.hide_gridlines(2)
        sheet.set_default_row(18)
        sheet.freeze_panes(4, 0)

    metadata = pack["metadata"]
    kpi = pack["kpis"]
    _write_title(
        dashboard,
        formats,
        "UniSweet Finance Analysis Dashboard",
        f"{metadata['reporting_month'][:7]} vs {metadata['comparison_month'][:7]} | {metadata['comparison_basis']} | {metadata['publication_status']}",
    )
    dashboard.set_column("A:O", 10)
    dashboard.set_column("D:D", 2)
    dashboard.set_column("H:H", 2)
    dashboard.set_column("L:L", 2)
    dashboard.merge_range("A5:O5", "EXECUTIVE SCORECARD — kEUR unless stated otherwise", formats["section"])
    cards = [
        ("A7:C7", "A8:C10", "TURNOVER", kpi["current_turnover_keur"], "card_blue"),
        ("E7:G7", "E8:G10", "TURNOVER VARIANCE", kpi["turnover_variance_keur"], "card_red"),
        ("I7:K7", "I8:K10", "DISCOUNT % TO", kpi["current_discount_pct_to"], "card_pct"),
        ("M7:O7", "M8:O10", "PBO", kpi["current_pbo_keur"], "card_green"),
    ]
    for label_range, value_range, label, value, style in cards:
        dashboard.merge_range(label_range, label, formats["card_label"])
        dashboard.merge_range(value_range, value, formats[style])
    dashboard.merge_range("A12:H12", "LARGEST NEGATIVE TURNOVER DRIVERS", formats["section"])
    dashboard.merge_range("I12:O12", "THREE ACTIONS FOR MANAGEMENT", formats["section"])

    insight_ranges = [("I13:O16", "insight_blue"), ("I17:O21", "insight_red"), ("I22:O26", "insight_green")]
    for insight, (cell_range, style) in zip(pack["insights"], insight_ranges, strict=True):
        text = (
            f"{insight['rank']}. {insight['headline']}\n"
            f"Evidence: {insight['evidence']}\n"
            f"Action: {insight['recommended_action']}\n"
            f"Owner / timing: {insight['proposed_owner']} — {insight['timing']}"
        )
        dashboard.merge_range(cell_range, text, formats[style])

    dashboard.merge_range("A29:K29", "24-MONTH TURNOVER TREND BY BRAND", formats["section"])
    dashboard.merge_range("M29:O29", "DECISION CONTEXT", formats["section"])
    dashboard.merge_range("M30:O34", pack["insights"][0]["headline"], formats["body_wrap"])
    dashboard.merge_range("M35:O39", f"Leakage proxy: {_signed(pack['insights'][1]['financial_impact_keur'])} kEUR", formats["body_wrap"])
    dashboard.merge_range("M40:O44", pack["insights"][2]["headline"], formats["body_wrap"])
    dashboard.merge_range("A46:O46", f"Input hash: {metadata['input_hash']} | Generated: {metadata['generated_at_utc']}", formats["section"])

    # Sales sheet and chart-driving tables.
    _write_title(sales, formats, "Sales performance and drivers", f"{metadata['comparison_basis']} comparison | calculations performed in Python", "T")
    sales.set_column("A:A", 14)
    sales.set_column("B:B", 15)
    sales.set_column("C:C", 24)
    sales.set_column("D:G", 16)
    sales.set_column("I:T", 14)
    scorecard_rows = [
        ["Turnover", kpi["current_turnover_keur"], kpi["prior_turnover_keur"], kpi["turnover_variance_keur"], kpi["turnover_growth_pct"]],
        ["GSV", kpi["current_gsv_keur"], kpi["prior_gsv_keur"], kpi["gsv_variance_keur"], kpi["gsv_growth_pct"]],
        ["Discount", kpi["current_discount_keur"], kpi["prior_discount_keur"], kpi["discount_variance_keur"], kpi["discount_growth_pct"]],
        ["Discount % TO", kpi["current_discount_pct_to"], kpi["prior_discount_pct_to"], None, kpi["discount_pct_to_movement_bps"]],
    ]
    _write_table(sales, 4, 0, ["Metric", "Current", "Prior", "Variance", "Growth / movement"], scorecard_rows, formats, {1: "amount", 2: "amount", 3: "amount", 4: "pct"})
    sales.write(8, 1, kpi["current_discount_pct_to"], formats["pct"])
    sales.write(8, 2, kpi["prior_discount_pct_to"], formats["pct"])
    sales.write(8, 4, kpi["discount_pct_to_movement_bps"], formats["bps"])

    driver_rows = [
        [row["analysis_dimension"], row["dimension_id"], row["dimension_name"], row["current_turnover_keur"], row["turnover_variance_keur"], row["turnover_growth_pct"], row["variance_contribution_pct"]]
        for row in pack["drivers"]
    ]
    _write_table(sales, 10, 0, ["Dimension", "ID", "Entity", "Current TO", "TO variance", "Growth", "Variance contribution"], driver_rows, formats, {3: "amount", 4: "amount", 5: "pct", 6: "pct"})
    trend_rows = [[row["reporting_month"][:7], row["gsv_keur"], row["discount_keur"], row["turnover_keur"]] for row in pack["trend"]]
    _write_table(sales, 10, 8, ["Month", "GSV", "Discount", "Turnover"], trend_rows, formats, {1: "amount", 2: "amount", 3: "amount"})

    brand_names = sorted({row["brand_name"] for row in pack["brand_trend"]})
    months = sorted({row["reporting_month"][:7] for row in pack["brand_trend"]})
    brand_values = {(row["reporting_month"][:7], row["brand_name"]): row["turnover_keur"] for row in pack["brand_trend"]}
    brand_rows = [[month, *[brand_values.get((month, brand)) for brand in brand_names]] for month in months]
    _write_table(sales, 10, 13, ["Month", *brand_names], brand_rows, formats, {column: "amount" for column in range(1, len(brand_names) + 1)})

    chart_driver_rows = sorted(pack["drivers"], key=lambda row: row["turnover_variance_keur"])[:7]
    _write_table(
        sales,
        10,
        18,
        ["Driver", "TO variance"],
        [[f"{row['analysis_dimension'][:3]} | {row['dimension_name']}", row["turnover_variance_keur"]] for row in chart_driver_rows],
        formats,
        {1: "amount"},
    )

    driver_chart = workbook.add_chart({"type": "column"})
    driver_chart.add_series(
        {
            "name": "TO variance",
            "categories": ["Sales Analysis", 11, 18, 10 + len(chart_driver_rows), 18],
            "values": ["Sales Analysis", 11, 19, 10 + len(chart_driver_rows), 19],
            "fill": {"color": "#247BA0"},
            "border": {"none": True},
        }
    )
    driver_chart.set_title({"name": "Largest negative Turnover drivers (kEUR)"})
    driver_chart.set_legend({"none": True})
    driver_chart.set_y_axis({"num_format": "#,##0;[Red](#,##0)", "major_gridlines": {"visible": True, "line": {"color": GRID}}})
    driver_chart.set_style(10)
    dashboard.insert_chart("A13", driver_chart, {"x_scale": 0.90, "y_scale": 0.90})

    brand_chart = workbook.add_chart({"type": "line"})
    for index, brand in enumerate(brand_names, start=1):
        brand_chart.add_series(
            {
                "name": ["Sales Analysis", 10, 13 + index],
                "categories": ["Sales Analysis", 11, 13, 10 + len(brand_rows), 13],
                "values": ["Sales Analysis", 11, 13 + index, 10 + len(brand_rows), 13 + index],
                "line": {"width": 2},
            }
        )
    brand_chart.set_title({"name": "Monthly Turnover by Brand (kEUR)"})
    brand_chart.set_legend({"position": "bottom"})
    brand_chart.set_y_axis({"num_format": "#,##0", "major_gridlines": {"visible": True, "line": {"color": GRID}}})
    brand_chart.set_style(10)
    dashboard.insert_chart("A30", brand_chart, {"x_scale": 1.40, "y_scale": 0.85})

    # Leakage sheet.
    _write_title(leakage, formats, "Discount leakage diagnostics", "Prior-period rate applied to current Turnover; diagnostic proxy only", "K")
    leakage.set_column("A:B", 15)
    leakage.set_column("C:C", 25)
    leakage.set_column("D:I", 16)
    leakage.set_column("J:J", 26)
    leakage_rows = [
        [row["analysis_dimension"], row["dimension_id"], row["dimension_name"], row["turnover_keur_current"], row["turnover_variance_keur"], row["current_discount_pct_to"], row["prior_discount_pct_to"], row["discount_pct_to_movement_bps"], row["discount_improvement_opportunity_keur"], row["leakage_quadrant"]]
        for row in pack["discount_leakage"]
    ]
    _write_table(leakage, 4, 0, ["Dimension", "ID", "Entity", "Current TO", "TO variance", "Current Disc % TO", "Prior Disc % TO", "Movement", "Recovery proxy", "Quadrant"], leakage_rows, formats, {3: "amount", 4: "amount", 5: "pct", 6: "pct", 7: "bps", 8: "amount"})
    leakage.conditional_format(5, 8, 4 + len(leakage_rows), 8, {"type": "data_bar", "bar_color": "#F28E8E"})

    # P&L and Market sheet.
    _write_title(pnl_market, formats, "P&L quality and Market context", "Annual P&L is certified; MAT Market comparison is directional", "N")
    pnl_market.set_column("A:A", 18)
    pnl_market.set_column("B:K", 16)
    pnl_market.set_column("L:N", 25)
    pnl_rows = [
        [row["brand_name"], row["turnover_keur_current"], row["turnover_variance_keur"], row["gross_profit_keur_current"], row["gross_profit_variance_keur"], row["pbo_keur_current"], row["pbo_variance_keur"], row["current_gross_margin_pct"], row["gross_margin_movement_bps"], row["current_pbo_pct_to"], row["pbo_margin_movement_bps"]]
        for row in sorted(pack["pnl"], key=lambda row: (row["brand_name"] != "TOTAL", row["brand_name"]))
    ]
    _write_table(pnl_market, 4, 0, ["Brand", "Current TO", "TO variance", "Current GP", "GP variance", "Current PBO", "PBO variance", "GM %", "GM movement", "PBO % TO", "PBO movement"], pnl_rows, formats, {1: "amount", 2: "amount", 3: "amount", 4: "amount", 5: "amount", 6: "amount", 7: "pct", 8: "bps", 9: "pct", 10: "bps"})
    market_rows = [
        [row["brand_name"], row["internal_calendar_year"], row["internal_turnover_growth_pct"], row["reporting_period"], row["market_value_growth_pct"], row["value_share_mat"], row["share_movement_pp"], row["sellin_vs_market_growth_gap_pp"], "CY sell-in vs supplied MAT market; directional only"]
        for row in pack["market"]
    ]
    _write_table(pnl_market, 12, 0, ["Brand", "Internal year", "Internal growth", "Market period", "Market growth", "Value share", "Share movement pp", "Sell-in vs Market gap pp", "Caveat"], market_rows, formats, {2: "pct", 4: "pct", 5: "pct", 8: "body_wrap"})
    pnl_market.set_column("I:I", 25)
    for row_index in range(13, 13 + len(market_rows)):
        pnl_market.set_row(row_index, 30)

    # Checks and source lineage.
    _write_title(checks, formats, "Checks and source lineage", f"MODEL STATUS: {metadata['publication_status']} | {metadata['critical_failure_count']} critical failures", "H")
    checks.set_column("A:A", 31)
    checks.set_column("B:B", 15)
    checks.set_column("C:C", 48)
    checks.set_column("D:F", 14)
    checks.set_column("G:H", 60)
    check_rows = [[row["check_id"], row["area"], row["name"], row["severity"], row["status"], row["exceptions"], row["details"]] for row in pack["checks"]]
    _write_table(checks, 4, 0, ["Check ID", "Area", "Check", "Severity", "Status", "Exceptions", "Details"], check_rows, formats, {6: "body_wrap"})
    for index, row in enumerate(pack["checks"], start=5):
        checks.write(index, 4, row["status"], formats[row["status"].lower()])
    source_start = 7 + len(check_rows)
    checks.merge_range(source_start, 0, source_start, 7, "SOURCE MANIFEST", formats["section"])
    source_rows = [[row["source_type"], row["path"], row["rows"], f"{row['sha256'][:12]}…", row["modified_at_utc"][:16].replace("T", " "), row["schema_status"], row["note"]] for row in pack["source_manifest"]]
    _write_table(checks, source_start + 1, 0, ["Type", "Path", "Rows", "SHA-256", "Modified UTC", "Schema", "Note"], source_rows, formats, {1: "body_wrap", 3: "body_wrap", 4: "body_wrap", 6: "body_wrap"})

    dashboard.set_landscape()
    dashboard.fit_to_pages(1, 1)
    dashboard.print_area("A1:O46")
    dashboard.set_zoom(85)
    workbook.close()


def write_story_review(pack: dict[str, Any], story: dict[str, Any], output_path: Path) -> None:
    lines = [
        "# Monthly Finance Story Review",
        "",
        f"- Story status: **{story['status'].upper()}**",
        f"- Reporting month: {pack['metadata']['reporting_month']}",
        f"- Comparison: {pack['metadata']['comparison_basis']} vs {pack['metadata']['comparison_month']}",
        f"- Data status: **{pack['metadata']['publication_status']}**",
        "",
        "> PowerPoint is generated only when `config/story.json` has `status: approved`.",
        "",
    ]
    for insight in pack["insights"]:
        lines.extend(
            [
                f"## {insight['rank']}. {insight['headline']}",
                "",
                f"**Evidence:** {insight['evidence']}",
                "",
                f"**Financial impact:** {_signed(insight['financial_impact_keur'])} kEUR — {insight['impact_basis']}",
                "",
                f"**Recommended action:** {insight['recommended_action']}",
                "",
                f"**Proposed owner / timing:** {insight['proposed_owner']} — {insight['timing']}",
                "",
                f"**Caveat:** {insight['caveat']}",
                "",
            ]
        )
    output_path.write_text("\n".join(lines), encoding="utf-8")


def _rgb(hex_color: str) -> RGBColor:
    value = hex_color.removeprefix("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _add_textbox(slide: Any, left: float, top: float, width: float, height: float, text: str, *, size: int, color: str = INK, bold: bool = False, fill: str | None = None, align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = _rgb(fill)
        shape.line.color.rgb = _rgb(fill)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    paragraph = frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    paragraph.font.name = "Segoe UI"
    paragraph.font.size = Pt(size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = _rgb(color)


def publish_presentation(pack: dict[str, Any], story: dict[str, Any], output_path: Path) -> None:
    if story.get("status", "draft").lower() != "approved":
        raise ValueError("Presentation publishing requires config/story.json status=approved")
    if pack["metadata"]["publication_status"] != "READY":
        raise ValueError("Presentation publishing requires publication_status=READY")
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = _rgb(WHITE)

    _add_textbox(slide, 0, 0, 13.333, 0.75, "Monthly Finance Review — decisions, not data theatre", size=26, color=WHITE, bold=True, fill=NAVY)
    _add_textbox(slide, 0.35, 0.78, 12.6, 0.32, f"{pack['metadata']['reporting_month'][:7]} | {pack['metadata']['comparison_basis']} | kEUR | input hash {pack['metadata']['input_hash'][:8]}", size=10, color="#64748B")

    kpi = pack["kpis"]
    cards = [
        ("TURNOVER", f"{kpi['current_turnover_keur']:,.0f}", BLUE),
        ("TO VARIANCE", f"{kpi['turnover_variance_keur']:,.0f}", RED),
        ("DISCOUNT % TO", _percent(kpi["current_discount_pct_to"]), GREEN),
        ("PBO", f"{kpi['current_pbo_keur']:,.0f}", GREEN),
    ]
    for index, (label, value, fill) in enumerate(cards):
        left = 0.35 + index * 3.2
        _add_textbox(slide, left, 1.15, 2.85, 0.28, label, size=10, color=WHITE, bold=True, fill=NAVY, align=PP_ALIGN.CENTER)
        _add_textbox(slide, left, 1.43, 2.85, 0.63, value, size=22, bold=True, fill=fill, align=PP_ALIGN.CENTER)

    _add_textbox(slide, 0.35, 2.25, 4.0, 0.35, "Largest Turnover declines", size=14, color=WHITE, bold=True, fill=NAVY)
    negative = sorted(
        [row for row in pack["drivers"] if row["turnover_variance_keur"] < 0],
        key=lambda row: row["turnover_variance_keur"],
    )[:6]
    maximum = max(abs(row["turnover_variance_keur"]) for row in negative) if negative else 1
    for index, row in enumerate(negative):
        top = 2.72 + index * 0.58
        label = f"{row['analysis_dimension'][:3]} | {row['dimension_name']}"
        _add_textbox(slide, 0.35, top, 1.45, 0.28, label, size=9)
        width = 2.15 * abs(row["turnover_variance_keur"]) / maximum
        _add_textbox(slide, 1.83, top, width, 0.28, f"{row['turnover_variance_keur']:,.0f}", size=8, color=WHITE, bold=True, fill="#247BA0", align=PP_ALIGN.RIGHT)

    _add_textbox(slide, 4.65, 2.25, 8.33, 0.35, "Three actions for management", size=14, color=WHITE, bold=True, fill=NAVY)
    fills = [BLUE, RED, GREEN]
    for index, insight in enumerate(pack["insights"]):
        top = 2.72 + index * 1.38
        text = (
            f"{insight['rank']}. {insight['headline']}\n"
            f"ACTION: {insight['recommended_action']}\n"
            f"OWNER / TIMING: {insight['proposed_owner']} — {insight['timing']}"
        )
        _add_textbox(slide, 4.65, top, 8.33, 1.15, text, size=10, fill=fills[index])

    _add_textbox(slide, 0.35, 6.83, 12.63, 0.35, "Caveat: discount opportunity is diagnostic; P&L is annual; Market comparison is MAT and directional.", size=9, color="#64748B")
    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(output_path)


def write_json(data: dict[str, Any], path: Path) -> None:
    path.write_text(json.dumps(data, indent=2, ensure_ascii=False) + "\n", encoding="utf-8")
