"""Slide-design driver: one spec, two renderings.

The slideument problem is that a live deck and a circulated deck want different
densities, and one artifact aimed at the middle serves neither. The answer here
is not to choose a density but to write the content once and render it twice:

    build_pptx(slides, path)   sparse - title + visual, everything else in notes
    build_pdf(slides, path)    annotated - takeaway on the page, narration under it

Both consume the same list of `Slide`, so the two files cannot drift apart.

Geometry is deliberately identical to the chart canvas. `visualisations.ipynb`
exports headerless twins to outputs/report_visuals/slides/ by hiding its
figure-level header and source line while leaving the axes untouched, so those
PNGs carry empty bands exactly where this module puts native text back. Nothing
is scaled or cropped; the charts keep the alignment they were tuned for.

Palette comes from the sibling storytelling-with-data skill. It is imported, not
copied, so the deck and the charts have one definition of the visual system.
"""

from __future__ import annotations

import sys
import textwrap
from dataclasses import dataclass, field
from pathlib import Path

HERE = Path(__file__).resolve().parent
# The chart toolkit is a sibling directory, but it is named differently depending on
# where this file lives: `charts` inside the repository, `storytelling-with-data`
# when the two are installed side by side as skills. Try both rather than pinning
# one layout and breaking the other.
for _candidate in (HERE.parent / "charts", HERE.parent / "storytelling-with-data"):
    if (_candidate / "swd.py").exists():
        if str(_candidate) not in sys.path:
            sys.path.insert(0, str(_candidate))
        break
else:  # pragma: no cover - a layout we have not seen
    raise ImportError(
        f"deck.py needs swd.py from the chart toolkit; looked for a sibling "
        f"'charts' or 'storytelling-with-data' directory next to {HERE}"
    )

import swd  # noqa: E402  - the one definition of the palette

__all__ = ["Slide", "build_pptx", "build_pdf", "W_IN", "H_IN"]

# --- geometry ----------------------------------------------------------------
# 16:9 at the charts' own export size. Fractions match add_header/add_footer.
W_IN, H_IN = 40 / 3, 7.5
LEFT = 0.055                      # the one left edge, for every rendering
RIGHT = 0.945
EYEBROW_Y = 0.980                 # the signpost: where the reader is
TITLE_Y = 0.940                   # top of the title block
SUB_Y = 0.886
TAKEAWAY_Y = 0.836                # circulated rendering only
NARRATION_Y = 0.088               # circulated rendering only
SOURCE_Y = 0.030

# In the circulated rendering the chart is cropped to its ink and placed in the
# band between the header and the narration, so the narration gets real space
# instead of landing on top of the axis labels. The bottom edge sits where it
# does because the narration below it is set at NARRATION_PT, not because 0.20
# is a round number - raise the type and this has to come down with it.
ART_TOP = 0.79
ART_BOTTOM = 0.20

TITLE_PT = 23.0
SUB_PT = 11.5
TAKEAWAY_PT = 13.0
NARRATION_PT = 12.0
SOURCE_PT = 8.5
FONT = "Arial"

# The one big line on a statement slide, and the character counts that keep it
# and its follow-on lines inside the right margin. The band is the space left
# between the takeaway above and the narration below; the block is centred in it
# because a statement slide has no artwork to hold the lower half of the page.
STATEMENT_PT = 26
STATEMENT_WRAP = 58
STATEMENT_BODY_WRAP = 132
STATEMENT_GAP = 0.055
STATEMENT_BAND_TOP = 0.78

# Label-and-body lists. MIN_ITEM_GAP is larger than LABEL_OFFSET on purpose:
# the space between two items has to beat the space inside one, or a label reads
# as belonging to the paragraph above it.
LABEL_OFFSET = 0.038
BODY_LINE_H = 0.035
MIN_ITEM_GAP = 0.048

# --- palette, inherited ------------------------------------------------------
INK = swd.INK
MID = swd.GREY
MUTED = swd.MUTED
ACCENT = swd.ACCENT
CONTRA = swd.ACCENT_2
PALE = swd.SEQ_BLUE[-1]

# Narration has to fit the band the chart leaves free at the bottom. Keeping it
# tight is also the right editorial discipline - crisp beats complete.
# The wrap is a character count standing in for a width, so it is tied to the
# point size: set larger type and the same count runs past the right margin.
NARRATION_WRAP = round(118 * 9.5 / NARRATION_PT)
NARRATION_MAX_LINES = 3
NARRATION_LINE_STEP = 0.028 * NARRATION_PT / 9.5


def fit_title(title: str) -> float:
    """Titles are edited to fit; rendering never shrinks the argument."""
    return TITLE_PT


def ink_bbox(png: Path, pad: int = 8) -> tuple[int, int, int, int]:
    """Crop box around everything that is not page white.

    The chart twins keep the full 16:9 canvas with empty header and footer
    bands. Full bleed is right for the .pptx, where native text fills those
    bands. The circulated page needs the bands back as narration space, so it
    crops to the ink instead of guessing fixed margins per chart.
    """
    import numpy as np
    from PIL import Image

    with Image.open(png) as im:
        arr = np.asarray(im.convert("L"))
    rows = np.where((arr < 250).any(axis=1))[0]
    cols = np.where((arr < 250).any(axis=0))[0]
    if not len(rows) or not len(cols):
        return (0, 0, arr.shape[1], arr.shape[0])
    return (max(int(cols[0]) - pad, 0), max(int(rows[0]) - pad, 0),
            min(int(cols[-1]) + pad, arr.shape[1]), min(int(rows[-1]) + pad, arr.shape[0]))


def wrap_narration(text: str) -> list[str]:
    lines = textwrap.wrap(text, NARRATION_WRAP) if text else []
    if len(lines) > NARRATION_MAX_LINES:
        raise ValueError(
            f"narration needs {len(lines)} lines, band fits {NARRATION_MAX_LINES}: {text[:70]}..."
        )
    return lines


def narration_top(text: str) -> float:
    """Figure fraction of the top of the narration block.

    Body copy has to stop above this. It is derived from the type size and the
    line count rather than being a constant, because the last constant here was
    chosen for 9.5pt narration and silently started overlapping the appendix
    body when the narration was set larger.
    """
    lines = wrap_narration(text)
    if not lines:
        return SOURCE_Y + 0.03
    line_h = NARRATION_PT * 1.45 / 72 / H_IN
    return NARRATION_Y + NARRATION_LINE_STEP * (len(lines) - 1) + len(lines) * line_h


@dataclass
class Slide:
    """One unit of the argument, in whichever rendering."""

    kind: str = "chart"          # title | statement | summary | section | chart | actions | closing
    title: str = ""              # the action title - the claim, never the topic
    subtitle: str = ""
    takeaway: str = ""           # the "so what", on the page when circulated
    narration: str = ""          # what the presenter would say
    image: Path | None = None    # headerless chart twin
    source: str = ""
    bullets: list = field(default_factory=list)   # str, or (label, body) pairs
    eyebrow: str = ""            # signpost: where the reader is

    def notes(self) -> str:
        parts = [p for p in (self.takeaway, self.narration) if p]
        if self.source:
            parts.append(f"[Sources]\n- {self.source}")
        return "\n\n".join(parts)


# --- rendering 1: the room ---------------------------------------------------
def build_pptx(slides: list[Slide], path: Path) -> Path:
    """Sparse slides. Takeaway and narration live in the speaker notes."""
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Inches, Pt

    def rgb(hexstr: str) -> RGBColor:
        return RGBColor.from_string(hexstr.lstrip("#").upper())

    def textbox(slide, x, y, w, h, text, size, colour, bold=False, spacing=1.15):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        for i, chunk in enumerate(text.split("\n")):
            para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            para.line_spacing = spacing
            run = para.add_run()          # never tf.text - it drops the styling
            run.text = chunk
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.name = FONT
            run.font.color.rgb = rgb(colour)
        return box

    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    blank = prs.slide_layouts[6]

    for spec in slides:
        slide = prs.slides.add_slide(blank)

        if spec.image:
            # Full bleed. The empty bands in the twin are where the text goes.
            slide.shapes.add_picture(str(spec.image), 0, 0,
                                     width=Inches(W_IN), height=Inches(H_IN))

        x = LEFT * W_IN
        width = (RIGHT - LEFT) * W_IN

        if spec.eyebrow:
            textbox(slide, x, (1 - EYEBROW_Y) * H_IN, width, 0.25,
                    spec.eyebrow.upper(), 9.5, MUTED, bold=True)

        if spec.kind == "title":
            textbox(slide, x, H_IN * 0.34, width, 1.6, spec.title, 40, ACCENT, bold=True)
            if spec.subtitle:
                textbox(slide, x, H_IN * 0.56, width, 1.0, spec.subtitle, 15, MID)
        elif spec.title:
            textbox(slide, x, (1 - TITLE_Y) * H_IN, width, 0.9,
                    spec.title, fit_title(spec.title), ACCENT, bold=True)
            if spec.subtitle:
                textbox(slide, x, (1 - SUB_Y) * H_IN, width, 0.6, spec.subtitle, SUB_PT, MID)

        if spec.bullets:
            _pptx_bullets(spec, slide, textbox, x, width)

        if spec.source:
            textbox(slide, x, H_IN - (SOURCE_Y * H_IN) - 0.20, width, 0.20,
                    f"Source: {spec.source}", SOURCE_PT, MUTED)

        # Read notes_slide before writing or python-pptx never creates the part.
        slide.notes_slide.notes_text_frame.text = spec.notes()

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path


def _pptx_bullets(spec, slide, textbox, x, width):
    """Statement, summary, actions and closing bodies."""
    top = H_IN * (0.26 if spec.kind == "statement" else 0.24)
    if spec.kind == "statement":
        # Same fix as the PDF: advance by the headline's own height rather than a
        # fixed 2.6 inches, which left most of a page empty under a single line,
        # and centre the block since there is no artwork holding the lower half.
        headline = textwrap.fill(spec.bullets[0], STATEMENT_WRAP)
        rows = headline.count("\n") + 1
        head_h = rows * (STATEMENT_PT * 1.30 / 72)
        body_h = sum(0.14 if not line else (0.55 if len(line) > 110 else 0.36)
                     for line in spec.bullets[1:])
        band_top, band_bottom = H_IN * 0.24, H_IN * 0.82
        top = band_top + max((band_bottom - band_top - head_h - 0.42 - body_h) / 2, 0.0)
        textbox(slide, x, top, width, rows * 0.62 + 0.2, headline,
                STATEMENT_PT, INK, bold=True, spacing=1.30)
        y = top + head_h + 0.42
        for line in spec.bullets[1:]:
            if not line:
                y += 0.14
                continue
            textbox(slide, x, y, width, 0.80, line, 13, MID, spacing=1.45)
            y += 0.55 if len(line) > 110 else 0.36
        return

    if spec.kind == "actions":
        col_w = (width - 0.7) / 3
        for i, (head, body) in enumerate(spec.bullets):
            cx = x + i * (col_w + 0.35)
            textbox(slide, cx, top, col_w, 0.9, head, 15, ACCENT, bold=True, spacing=1.20)
            textbox(slide, cx, top + 1.0, col_w, 3.0, body, 12, INK, spacing=1.40)
        return

    for i, item in enumerate(spec.bullets):
        y = top + i * 0.92
        if isinstance(item, (tuple, list)):
            label, body = item
            textbox(slide, x, y, width, 0.28, label, 12, ACCENT if spec.kind == "closing" else MID, bold=True)
            textbox(slide, x, y + 0.28, width, 0.52, body, 14, INK, spacing=1.25)
        else:
            textbox(slide, x, y, width, 0.5, item, 15, INK, spacing=1.30)


# --- rendering 2: the reader -------------------------------------------------
def build_pdf(slides: list[Slide], path: Path) -> Path:
    """Same title and visual, plus the takeaway and narration on the page."""
    import matplotlib
    matplotlib.use("Agg")
    import matplotlib.image as mpimg
    import matplotlib.pyplot as plt
    from matplotlib.backends.backend_pdf import PdfPages

    plt.rcParams.update({
        "font.family": "sans-serif",
        "font.sans-serif": [FONT, "Helvetica", "DejaVu Sans"],
        "pdf.fonttype": 42,          # embed as TrueType so the text stays selectable
    })

    path.parent.mkdir(parents=True, exist_ok=True)
    with PdfPages(str(path)) as pdf:
        for spec in slides:
            fig = plt.figure(figsize=(W_IN, H_IN), dpi=144)
            fig.patch.set_facecolor("white")

            if spec.image:
                _place_art(fig, spec.image)

            if spec.eyebrow:
                fig.text(LEFT, EYEBROW_Y, spec.eyebrow.upper(), ha="left", va="top",
                         fontsize=9.5, color=MUTED, fontweight="bold")

            if spec.kind == "title":
                fig.text(LEFT, 0.62, spec.title, ha="left", va="top",
                         fontsize=40, color=ACCENT, fontweight="bold")
                if spec.subtitle:
                    fig.text(LEFT, 0.44, spec.subtitle, ha="left", va="top",
                             fontsize=15, color=MID, linespacing=1.5)
            elif spec.title:
                fig.text(LEFT, TITLE_Y, spec.title, ha="left", va="top",
                         fontsize=fit_title(spec.title), color=ACCENT, fontweight="bold")
                if spec.subtitle:
                    fig.text(LEFT, SUB_Y, spec.subtitle, ha="left", va="top",
                             fontsize=SUB_PT, color=MID)

            # The annotation layer - this is what replaces the presenter.
            if spec.takeaway:
                fig.text(LEFT, TAKEAWAY_Y, spec.takeaway, ha="left", va="top",
                         fontsize=TAKEAWAY_PT, color=INK,
                         fontweight="bold", linespacing=1.35)

            if spec.bullets:
                _pdf_bullets(spec, fig)

            # The closing line of every page, and part of the annotation layer:
            # it carries the same blue as the callouts inside the charts, so
            # on-page annotation reads as one thing whether it sits in the plot
            # or beneath it. The takeaway keeps the hierarchy by being bold INK.
            lines = wrap_narration(spec.narration)
            if lines:
                fig.text(LEFT, NARRATION_Y + NARRATION_LINE_STEP * (len(lines) - 1),
                         "\n".join(lines), ha="left", va="bottom",
                         fontsize=NARRATION_PT, color=ACCENT, linespacing=1.45)

            if spec.source:
                fig.text(LEFT, SOURCE_Y, f"Source: {spec.source}", ha="left", va="bottom",
                         fontsize=SOURCE_PT, color=MUTED)

            pdf.savefig(fig, facecolor="white")
            plt.close(fig)
    return path


def _place_art(fig, image: Path) -> None:
    """Crop the chart to its ink and centre it in the band left for artwork.

    Aspect is preserved, so a wide chart keeps its width and a tall one keeps
    its height; whichever dimension binds decides the other.
    """
    import matplotlib.image as mpimg

    x0, y0, x1, y1 = ink_bbox(image)
    art = mpimg.imread(str(image))[y0:y1, x0:x1]

    box_w, box_h = RIGHT - LEFT, ART_TOP - ART_BOTTOM
    page_aspect = W_IN / H_IN
    art_aspect = (art.shape[1] / art.shape[0]) / page_aspect   # in figure-fraction terms

    if art_aspect >= box_w / box_h:      # width-bound
        w, h = box_w, box_w / art_aspect
    else:                                 # height-bound
        w, h = box_h * art_aspect, box_h

    left = LEFT + (box_w - w) / 2
    bottom = ART_BOTTOM + (box_h - h) / 2
    ax = fig.add_axes([left, bottom, w, h])
    ax.imshow(art, interpolation="antialiased")
    ax.axis("off")


def _pdf_bullets(spec, fig) -> None:
    # A slide with no takeaway leaves the takeaway band empty; start the body
    # there rather than below a line that was never drawn.
    top = 0.74 if spec.kind == "statement" else (0.76 if spec.takeaway else 0.845)

    if spec.kind == "statement":
        # A statement slide carries no artwork, so the block is measured and then
        # centred in the band between the takeaway and the narration. The old code
        # pinned the headline to the top and jumped to a fixed y for the rest,
        # which left a quarter of the page empty in the middle of the argument and
        # the rest of it empty underneath.
        headline = textwrap.fill(spec.bullets[0], STATEMENT_WRAP)
        rows = headline.count("\n") + 1
        line_h = STATEMENT_PT * 1.30 / 72 / H_IN     # points -> figure fraction

        body = [textwrap.fill(line, STATEMENT_BODY_WRAP) if line else "" for line in spec.bullets[1:]]
        block = rows * line_h + STATEMENT_GAP + sum(
            0.030 if not w else 0.055 * (w.count("\n") + 1) + 0.020 for w in body
        )
        band_bottom = narration_top(spec.narration) + 0.030
        top = STATEMENT_BAND_TOP - max((STATEMENT_BAND_TOP - band_bottom - block) / 2, 0.0)

        fig.text(LEFT, top, headline, ha="left", va="top",
                 fontsize=STATEMENT_PT, color=INK, fontweight="bold", linespacing=1.30)

        y = top - rows * line_h - STATEMENT_GAP
        for wrapped in body:
            if not wrapped:                           # a blank entry is a spacer,
                y -= 0.030                            # not a whole line of space
                continue
            fig.text(LEFT, y, wrapped, ha="left", va="top",
                     fontsize=13, color=MID, linespacing=1.55)
            y -= 0.055 * (wrapped.count("\n") + 1) + 0.020
        return

    if spec.kind == "actions":
        for i, (head, body) in enumerate(spec.bullets):
            x = LEFT + i * 0.302
            fig.text(x, top, head, ha="left", va="top",
                     fontsize=15, color=ACCENT, fontweight="bold", linespacing=1.30)
            fig.text(x, top - 0.13, body, ha="left", va="top",
                     fontsize=12, color=INK, linespacing=1.55)
        return

    # Bodies wrap to the right margin rather than running off the page.
    BODY_WRAP = 116
    wrapped = [
        (item[0], textwrap.fill(item[1], BODY_WRAP)) if isinstance(item, (tuple, list))
        else textwrap.fill(item, BODY_WRAP)
        for item in spec.bullets
    ]

    # Space between items must exceed space inside an item, or the eye groups a
    # label with the body above it instead of the one below. So each item is
    # measured, and what is left over becomes the gap between them - rather than
    # one uniform step that shrinks the two distances together.
    heights = []
    for item in wrapped:
        body = item[1] if isinstance(item, tuple) else item
        rows = body.count("\n") + 1
        heights.append((LABEL_OFFSET if isinstance(item, tuple) else 0.0) + rows * BODY_LINE_H)

    # Stop above the narration, wherever it actually starts. This used to be a
    # flat 0.20, which was right only while the narration was set at 9.5pt.
    floor = narration_top(spec.narration) + 0.030
    available = top - floor
    gaps = max(len(wrapped) - 1, 1)
    gap = (available - sum(heights)) / gaps
    if gap < MIN_ITEM_GAP:
        raise ValueError(
            f"'{spec.title[:48]}' needs {sum(heights) + gaps * MIN_ITEM_GAP:.3f} of the page "
            f"but only {available:.3f} is free between the header and the narration. "
            f"Shorten the bullets or the narration."
        )

    y = top
    for item, height in zip(wrapped, heights):
        if isinstance(item, tuple):
            label, body = item
            fig.text(LEFT, y, label, ha="left", va="top", fontsize=12,
                     color=ACCENT if spec.kind == "closing" else MID, fontweight="bold")
            fig.text(LEFT, y - LABEL_OFFSET, body, ha="left", va="top", fontsize=14,
                     color=INK, linespacing=1.35)
        else:
            fig.text(LEFT, y, item, ha="left", va="top", fontsize=15,
                     color=INK, linespacing=1.40)
        y -= height + gap
