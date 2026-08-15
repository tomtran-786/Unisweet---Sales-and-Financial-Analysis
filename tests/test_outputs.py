from __future__ import annotations

from pathlib import Path

from pptx import Presentation

from unisweet_analysis.publish import publish_presentation


def test_analysis_pack_and_story_are_published_but_dashboard_and_draft_ppt_are_not(
    generated: tuple[dict, dict, Path],
) -> None:
    summary, _, output_dir = generated
    assert summary["story_status"] == "DRAFT"
    assert summary["presentation_status"] == "NOT_PUBLISHED"
    assert (output_dir / "analysis_pack.json").exists()
    assert (output_dir / "story_review.md").exists()
    assert not (output_dir / "finance_dashboard.xlsx").exists()
    assert not (output_dir / "monthly_review.pptx").exists()


def test_approved_story_can_publish_one_slide_deck(
    generated: tuple[dict, dict, Path], tmp_path: Path
) -> None:
    _, pack, _ = generated
    output_path = tmp_path / "approved_review.pptx"
    publish_presentation(pack, {"status": "approved"}, output_path)
    presentation = Presentation(output_path)
    assert len(presentation.slides) == 1
    text = "\n".join(shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text"))
    assert "Monthly Finance Review" in text
    assert "Three actions for management" in text
